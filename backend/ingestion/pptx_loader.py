from pptx import Presentation
from langchain_core.documents import Document

import os

def load_pptx_folder(folder_path, document_type="general"):
    documents = []

    for file in os.listdir(folder_path):
        if file.endswith(".pptx"):
            file_path = os.path.join(folder_path, file)
            prs = Presentation(file_path)

            for slide_number, slide in enumerate(prs.slides, start=1):
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)

                if slide_text:
                    documents.append(
                        Document(
                            page_content="\n".join(slide_text),
                            metadata={
                                "source": file,
                                "slide": slide_number,
                                "document_type": document_type
                            }
                        )
                    )

    return documents
