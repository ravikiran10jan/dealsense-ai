"""
Multimodal Document Processor for DealSense AI.

Uses GPT-4o Vision to describe images, charts, tables, and diagrams,
then stores the text descriptions in the existing TF-IDF/FAISS pipeline.

Approach: Vision-LLM + Text Embedding (Option A)
  Image/PDF/Slide -> GPT-4o Vision -> Text Description -> TF-IDF -> FAISS

Supported formats:
  - Images: PNG, JPG, JPEG, GIF, WEBP
  - PDFs: Extracts text + renders pages as images for visual content
  - PPTX: Extracts text + images from slides
"""
import os
import io
import base64
import logging
from typing import List, Dict, Any, Optional, Tuple

from langchain_core.documents import Document

logger = logging.getLogger(__name__)

# Supported file extensions
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp"}
DOCUMENT_EXTENSIONS = {".pdf"}
PRESENTATION_EXTENSIONS = {".pptx"}
ALL_SUPPORTED_EXTENSIONS = IMAGE_EXTENSIONS | DOCUMENT_EXTENSIONS | PRESENTATION_EXTENSIONS


def _get_vision_client():
    """Get OpenAI client for vision API calls."""
    from openai import OpenAI
    from dotenv import load_dotenv
    load_dotenv()
    return OpenAI(api_key=os.getenv("OPENAI_API_KEY"))


def _encode_image_to_base64(image_bytes: bytes) -> str:
    """Encode image bytes to base64 string."""
    return base64.b64encode(image_bytes).decode("utf-8")


def describe_image_with_vision(
    image_bytes: bytes,
    context: str = "",
    mime_type: str = "image/png",
) -> str:
    """
    Use GPT-4o Vision to describe an image.

    Args:
        image_bytes: Raw image bytes.
        context: Optional context about the image (e.g., slide title, source doc).
        mime_type: MIME type of the image.

    Returns:
        Text description of the image content.
    """
    client = _get_vision_client()
    b64_image = _encode_image_to_base64(image_bytes)

    system_prompt = (
        "You are an expert document analyst for a sales intelligence system. "
        "Describe the visual content in detail, focusing on: "
        "data in charts/tables (numbers, trends, comparisons), "
        "architecture diagrams (components, connections, flow), "
        "text content visible in screenshots, "
        "key takeaways and business insights. "
        "Be factual and structured. Output plain text, not markdown."
    )

    user_content = [
        {
            "type": "image_url",
            "image_url": {
                "url": f"data:{mime_type};base64,{b64_image}",
                "detail": "high",
            },
        },
    ]

    if context:
        user_content.insert(0, {
            "type": "text",
            "text": f"Context: {context}\n\nDescribe the visual content in this image:",
        })
    else:
        user_content.insert(0, {
            "type": "text",
            "text": "Describe the visual content in this image in detail:",
        })

    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_content},
            ],
            max_tokens=1000,
            temperature=0,
        )
        description = response.choices[0].message.content.strip()
        logger.info(f"Vision description generated ({len(description)} chars)")
        return description
    except Exception as e:
        logger.error(f"Vision API call failed: {e}")
        return f"[Image description unavailable: {e}]"


# ---------------------------------------------------------------------------
# Image file processing
# ---------------------------------------------------------------------------


def process_image_file(
    file_path: str,
    metadata: Optional[Dict[str, Any]] = None,
) -> List[Document]:
    """
    Process a standalone image file and return LangChain Documents.

    Args:
        file_path: Path to the image file.
        metadata: Additional metadata to attach.

    Returns:
        List of Document objects with the image description as page_content.
    """
    ext = os.path.splitext(file_path)[1].lower()
    mime_map = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".webp": "image/webp",
    }
    mime_type = mime_map.get(ext, "image/png")

    with open(file_path, "rb") as f:
        image_bytes = f.read()

    context = f"File: {os.path.basename(file_path)}"
    description = describe_image_with_vision(image_bytes, context=context, mime_type=mime_type)

    doc_metadata = {
        "source": os.path.basename(file_path),
        "type": "image",
        "content_type": "vision_description",
        "original_format": ext,
    }
    if metadata:
        doc_metadata.update(metadata)

    return [
        Document(
            page_content=f"[Image: {os.path.basename(file_path)}]\n\n{description}",
            metadata=doc_metadata,
        )
    ]


def process_image_bytes(
    image_bytes: bytes,
    filename: str = "uploaded_image",
    mime_type: str = "image/png",
    metadata: Optional[Dict[str, Any]] = None,
) -> List[Document]:
    """
    Process raw image bytes and return LangChain Documents.

    Args:
        image_bytes: Raw image data.
        filename: Name for the image.
        mime_type: MIME type.
        metadata: Additional metadata.

    Returns:
        List of Document objects.
    """
    description = describe_image_with_vision(
        image_bytes,
        context=f"File: {filename}",
        mime_type=mime_type,
    )

    doc_metadata = {
        "source": filename,
        "type": "image",
        "content_type": "vision_description",
    }
    if metadata:
        doc_metadata.update(metadata)

    return [
        Document(
            page_content=f"[Image: {filename}]\n\n{description}",
            metadata=doc_metadata,
        )
    ]


# ---------------------------------------------------------------------------
# PDF processing with image extraction
# ---------------------------------------------------------------------------


def process_pdf_file(
    file_path: str,
    metadata: Optional[Dict[str, Any]] = None,
    describe_pages: bool = True,
) -> List[Document]:
    """
    Process a PDF file: extract text and optionally describe pages visually.

    Uses PyMuPDF (fitz) for text extraction and page rendering.

    Args:
        file_path: Path to the PDF file.
        metadata: Additional metadata.
        describe_pages: If True, render pages with visual content and describe them.

    Returns:
        List of Document objects (text chunks + image descriptions).
    """
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    documents = []
    filename = os.path.basename(file_path)

    for page_num in range(len(doc)):
        page = doc[page_num]

        # 1. Extract text from the page
        text = page.get_text("text").strip()
        if text:
            doc_metadata = {
                "source": filename,
                "page": page_num + 1,
                "type": "pdf_text",
                "document_type": "pdf",
            }
            if metadata:
                doc_metadata.update(metadata)

            documents.append(
                Document(
                    page_content=f"[PDF: {filename}, Page {page_num + 1}]\n\n{text}",
                    metadata=doc_metadata,
                )
            )

        # 2. Check for images / visual content on the page
        if describe_pages:
            images = page.get_images(full=True)
            has_drawings = len(page.get_drawings()) > 5  # heuristic for charts/diagrams

            if images or has_drawings:
                # Render the page as an image and describe it
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes("png")

                context = f"PDF: {filename}, Page {page_num + 1}"
                if text:
                    # Give first 200 chars of text as context
                    context += f". Page text excerpt: {text[:200]}"

                description = describe_image_with_vision(
                    img_bytes,
                    context=context,
                    mime_type="image/png",
                )

                vis_metadata = {
                    "source": filename,
                    "page": page_num + 1,
                    "type": "pdf_visual",
                    "content_type": "vision_description",
                    "document_type": "pdf",
                    "image_count": len(images),
                }
                if metadata:
                    vis_metadata.update(metadata)

                documents.append(
                    Document(
                        page_content=(
                            f"[PDF Visual: {filename}, Page {page_num + 1}]\n\n{description}"
                        ),
                        metadata=vis_metadata,
                    )
                )

    doc.close()
    logger.info(f"Processed PDF '{filename}': {len(documents)} documents extracted")
    return documents


def process_pdf_bytes(
    pdf_bytes: bytes,
    filename: str = "uploaded.pdf",
    metadata: Optional[Dict[str, Any]] = None,
    describe_pages: bool = True,
) -> List[Document]:
    """
    Process PDF from raw bytes.

    Args:
        pdf_bytes: Raw PDF data.
        filename: Name for the document.
        metadata: Additional metadata.
        describe_pages: If True, describe pages with visual content.

    Returns:
        List of Document objects.
    """
    import fitz  # PyMuPDF

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    documents = []

    for page_num in range(len(doc)):
        page = doc[page_num]

        text = page.get_text("text").strip()
        if text:
            doc_metadata = {
                "source": filename,
                "page": page_num + 1,
                "type": "pdf_text",
                "document_type": "pdf",
            }
            if metadata:
                doc_metadata.update(metadata)

            documents.append(
                Document(
                    page_content=f"[PDF: {filename}, Page {page_num + 1}]\n\n{text}",
                    metadata=doc_metadata,
                )
            )

        if describe_pages:
            images = page.get_images(full=True)
            has_drawings = len(page.get_drawings()) > 5

            if images or has_drawings:
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes("png")

                context = f"PDF: {filename}, Page {page_num + 1}"
                if text:
                    context += f". Page text excerpt: {text[:200]}"

                description = describe_image_with_vision(
                    img_bytes, context=context, mime_type="image/png",
                )

                vis_metadata = {
                    "source": filename,
                    "page": page_num + 1,
                    "type": "pdf_visual",
                    "content_type": "vision_description",
                    "document_type": "pdf",
                    "image_count": len(images),
                }
                if metadata:
                    vis_metadata.update(metadata)

                documents.append(
                    Document(
                        page_content=(
                            f"[PDF Visual: {filename}, Page {page_num + 1}]\n\n{description}"
                        ),
                        metadata=vis_metadata,
                    )
                )

    doc.close()
    logger.info(f"Processed PDF '{filename}': {len(documents)} documents extracted")
    return documents


# ---------------------------------------------------------------------------
# PPTX processing with image extraction
# ---------------------------------------------------------------------------


def process_pptx_file(
    file_path: str,
    metadata: Optional[Dict[str, Any]] = None,
    describe_images: bool = True,
) -> List[Document]:
    """
    Process a PPTX file: extract text and describe embedded images.

    Args:
        file_path: Path to the PPTX file.
        metadata: Additional metadata.
        describe_images: If True, describe images found in slides.

    Returns:
        List of Document objects.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    documents = []
    filename = os.path.basename(file_path)

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        slide_images = []

        for shape in slide.shapes:
            # Extract text
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())

            # Extract images
            if describe_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    img_bytes = image.blob
                    content_type = image.content_type or "image/png"
                    slide_images.append((img_bytes, content_type))
                except Exception as e:
                    logger.warning(f"Could not extract image from slide {slide_number}: {e}")

        # Create text document for the slide
        if slide_texts:
            text_content = "\n".join(slide_texts)
            doc_metadata = {
                "source": filename,
                "slide": slide_number,
                "type": "pptx_text",
                "document_type": "presentation",
            }
            if metadata:
                doc_metadata.update(metadata)

            documents.append(
                Document(
                    page_content=text_content,
                    metadata=doc_metadata,
                )
            )

        # Describe each image on the slide
        for img_idx, (img_bytes, content_type) in enumerate(slide_images):
            slide_context = f"PPTX: {filename}, Slide {slide_number}"
            if slide_texts:
                slide_context += f". Slide text: {' '.join(slide_texts)[:200]}"

            description = describe_image_with_vision(
                img_bytes,
                context=slide_context,
                mime_type=content_type,
            )

            img_metadata = {
                "source": filename,
                "slide": slide_number,
                "image_index": img_idx + 1,
                "type": "pptx_image",
                "content_type": "vision_description",
                "document_type": "presentation",
            }
            if metadata:
                img_metadata.update(metadata)

            documents.append(
                Document(
                    page_content=(
                        f"[Slide {slide_number} Image {img_idx + 1}: {filename}]\n\n{description}"
                    ),
                    metadata=img_metadata,
                )
            )

    logger.info(f"Processed PPTX '{filename}': {len(documents)} documents extracted")
    return documents


def process_pptx_bytes(
    pptx_bytes: bytes,
    filename: str = "uploaded.pptx",
    metadata: Optional[Dict[str, Any]] = None,
    describe_images: bool = True,
) -> List[Document]:
    """
    Process PPTX from raw bytes.

    Args:
        pptx_bytes: Raw PPTX data.
        filename: Name for the document.
        metadata: Additional metadata.
        describe_images: If True, describe images found in slides.

    Returns:
        List of Document objects.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(pptx_bytes))
    documents = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        slide_images = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())

            if describe_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    img_bytes = image.blob
                    content_type = image.content_type or "image/png"
                    slide_images.append((img_bytes, content_type))
                except Exception as e:
                    logger.warning(f"Could not extract image from slide {slide_number}: {e}")

        if slide_texts:
            text_content = "\n".join(slide_texts)
            doc_metadata = {
                "source": filename,
                "slide": slide_number,
                "type": "pptx_text",
                "document_type": "presentation",
            }
            if metadata:
                doc_metadata.update(metadata)

            documents.append(
                Document(page_content=text_content, metadata=doc_metadata)
            )

        for img_idx, (img_bytes, content_type) in enumerate(slide_images):
            slide_context = f"PPTX: {filename}, Slide {slide_number}"
            if slide_texts:
                slide_context += f". Slide text: {' '.join(slide_texts)[:200]}"

            description = describe_image_with_vision(
                img_bytes, context=slide_context, mime_type=content_type,
            )

            img_metadata = {
                "source": filename,
                "slide": slide_number,
                "image_index": img_idx + 1,
                "type": "pptx_image",
                "content_type": "vision_description",
                "document_type": "presentation",
            }
            if metadata:
                img_metadata.update(metadata)

            documents.append(
                Document(
                    page_content=(
                        f"[Slide {slide_number} Image {img_idx + 1}: {filename}]\n\n{description}"
                    ),
                    metadata=img_metadata,
                )
            )

    logger.info(f"Processed PPTX '{filename}': {len(documents)} documents extracted")
    return documents


# ---------------------------------------------------------------------------
# Unified entry point
# ---------------------------------------------------------------------------


def process_document(
    file_path: str,
    metadata: Optional[Dict[str, Any]] = None,
) -> List[Document]:
    """
    Process any supported document file and return LangChain Documents.

    Dispatches to the appropriate processor based on file extension.

    Args:
        file_path: Path to the document file.
        metadata: Additional metadata.

    Returns:
        List of Document objects.

    Raises:
        ValueError: If the file extension is not supported.
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext in IMAGE_EXTENSIONS:
        return process_image_file(file_path, metadata=metadata)
    elif ext in DOCUMENT_EXTENSIONS:
        return process_pdf_file(file_path, metadata=metadata)
    elif ext in PRESENTATION_EXTENSIONS:
        return process_pptx_file(file_path, metadata=metadata)
    else:
        raise ValueError(
            f"Unsupported file format: '{ext}'. "
            f"Supported: {', '.join(sorted(ALL_SUPPORTED_EXTENSIONS))}"
        )


def process_document_bytes(
    file_bytes: bytes,
    filename: str,
    mime_type: Optional[str] = None,
    metadata: Optional[Dict[str, Any]] = None,
) -> List[Document]:
    """
    Process document from raw bytes.

    Dispatches to the appropriate processor based on filename extension or MIME type.

    Args:
        file_bytes: Raw file data.
        filename: Original filename (used to determine type).
        mime_type: Optional MIME type override.
        metadata: Additional metadata.

    Returns:
        List of Document objects.

    Raises:
        ValueError: If the file type is not supported.
    """
    ext = os.path.splitext(filename)[1].lower()

    if ext in IMAGE_EXTENSIONS:
        resolved_mime = mime_type or {
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".gif": "image/gif",
            ".webp": "image/webp",
        }.get(ext, "image/png")
        return process_image_bytes(
            file_bytes, filename=filename, mime_type=resolved_mime, metadata=metadata,
        )
    elif ext in DOCUMENT_EXTENSIONS:
        return process_pdf_bytes(file_bytes, filename=filename, metadata=metadata)
    elif ext in PRESENTATION_EXTENSIONS:
        return process_pptx_bytes(file_bytes, filename=filename, metadata=metadata)
    else:
        raise ValueError(
            f"Unsupported file format: '{ext}'. "
            f"Supported: {', '.join(sorted(ALL_SUPPORTED_EXTENSIONS))}"
        )
