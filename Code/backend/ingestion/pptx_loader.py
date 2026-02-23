from pptx import Presentation
from langchain_core.documents import Document

import os
import logging

logger = logging.getLogger(__name__)


def load_pptx_folder(folder_path, document_type="general", multimodal=False):
    """
    Load PPTX files from a folder and return LangChain Documents.

    Args:
        folder_path: Path to the folder containing PPTX files.
        document_type: Type label for metadata (e.g., 'case_study', 'offering').
        multimodal: If True, also extract and describe images using GPT-4o Vision.

    Returns:
        List of Document objects.
    """
    documents = []

    for file in os.listdir(folder_path):
        if file.endswith(".pptx"):
            file_path = os.path.join(folder_path, file)

            if multimodal:
                try:
                    from ingestion.multimodal_processor import process_pptx_file
                    docs = process_pptx_file(
                        file_path,
                        metadata={"document_type": document_type},
                    )
                    documents.extend(docs)
                    logger.info(f"Multimodal processed {file}: {len(docs)} documents")
                    continue
                except Exception as e:
                    logger.warning(f"Multimodal processing failed for {file}, falling back to text-only: {e}")

            # Text-only extraction (original behavior)
            prs = Presentation(file_path)

            for slide_number, slide in enumerate(prs.slides, start=1):
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)

                if slide_text:
                    documents.append(
                        Document(
                            page_content="\n".join(slide_text),
                            metadata={
                                "source": file,
                                "slide": slide_number,
                                "document_type": document_type
                            }
                        )
                    )

    return documents
